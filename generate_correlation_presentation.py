#!/usr/bin/env python3
from __future__ import annotations

import argparse
import random
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(221, 235, 247)
ACCENT = RGBColor(91, 155, 213)
DARK = RGBColor(30, 30, 30)


def add_header_band(slide):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.5))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()


def style_title(shape):
    p = shape.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(34)
    p.font.color.rgb = BLUE


def style_subtitle(shape):
    for p in shape.text_frame.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = DARK


def add_bullets(slide, title, bullets):
    if len(bullets) > 5:
        raise ValueError("Each slide can have at most 5 bullet points.")
    add_header_band(slide)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.65), Inches(6.2), Inches(0.8))
    title_box.text = title
    style_title(title_box)
    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.9), Inches(4.8))
    tf = body.text_frame
    tf.clear()
    for i, point in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = point
        p.level = 0
        p.space_after = Pt(12)
        p.font.size = Pt(22)
        p.font.color.rgb = DARK


def add_icon_circles(slide, labels):
    for i, label in enumerate(labels):
        x = Inches(7.2 + i * 1.8)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, Inches(5.5), Inches(1.2), Inches(1.2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = LIGHT_BLUE
        circle.line.color.rgb = ACCENT
        t = slide.shapes.add_textbox(x, Inches(5.88), Inches(1.2), Inches(0.5))
        p = t.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.alignment = 1


def draw_scatter_panel(slide, left, top, width, height, mode, title):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(248, 251, 255)
    panel.line.color.rgb = ACCENT

    plot_left = left + Inches(0.35)
    plot_top = top + Inches(0.5)
    plot_w = width - Inches(0.6)
    plot_h = height - Inches(0.9)

    x_axis = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, plot_left, plot_top + plot_h, plot_w, Inches(0.001))
    x_axis.line.color.rgb = BLUE
    y_axis = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, plot_left, plot_top, Inches(0.001), plot_h)
    y_axis.line.color.rgb = BLUE

    rng = random.Random(42)
    plot_w_in = plot_w / 914400.0
    plot_h_in = plot_h / 914400.0
    for i in range(22):
        x = i / 21
        jitter = rng.uniform(-0.09, 0.09)
        if mode == "positive":
            y = min(max(0.14 + 0.75 * x + jitter, 0.05), 0.95)
        elif mode == "negative":
            y = min(max(0.87 - 0.70 * x + jitter, 0.05), 0.95)
        elif mode == "zero":
            y = min(max(0.50 + rng.uniform(-0.34, 0.34), 0.05), 0.95)
        elif mode == "nonlinear":
            y = min(max(0.5 + 0.32 * (1 - (2 * x - 1) ** 2) + jitter * 0.6, 0.05), 0.95)
        else:
            y = min(max(0.10 + 0.80 * x + jitter * 0.8, 0.05), 0.95)
        dot_x = plot_left + Inches(x * (plot_w_in - 0.12))
        dot_y = plot_top + Inches((1 - y) * (plot_h_in - 0.12))
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, dot_x, dot_y, Inches(0.11), Inches(0.11))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()

    t = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.2), Inches(0.3))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLUE


def build_presentation(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_band(slide)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.2), Inches(1.2))
    title.text = "Measure of Correlation"
    style_title(title)
    subtitle = slide.shapes.add_textbox(Inches(0.82), Inches(2.5), Inches(8.5), Inches(0.9))
    subtitle.text = "Understanding Relationships Between Variables"
    style_subtitle(subtitle)
    draw_scatter_panel(slide, Inches(8.0), Inches(1.3), Inches(4.8), Inches(3.2), "positive", "Scatter Plot")
    add_icon_circles(slide, ["X", "Y", "r"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Introduction to Correlation",
        [
            "Correlation shows how two variables move together.",
            "Purpose: describe strength and direction of relationships.",
            "Types: positive, negative, and zero correlation.",
            "Example: study time vs. exam scores (positive).",
            "Example: product price vs. demand (negative).",
        ],
    )
    draw_scatter_panel(slide, Inches(6.9), Inches(1.1), Inches(2.0), Inches(2.0), "positive", "Positive")
    draw_scatter_panel(slide, Inches(9.0), Inches(1.1), Inches(2.0), Inches(2.0), "negative", "Negative")
    draw_scatter_panel(slide, Inches(11.1), Inches(1.1), Inches(2.0), Inches(2.0), "zero", "Zero")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Linear Correlation",
        [
            "A linear relationship follows a roughly straight pattern.",
            "Equal changes in X tend to produce proportional changes in Y.",
            "Scatter plots help reveal trend direction and consistency.",
            "Tight clustering means stronger linear correlation.",
            "Curved patterns indicate nonlinearity, not linear fit.",
        ],
    )
    draw_scatter_panel(slide, Inches(6.9), Inches(1.2), Inches(2.9), Inches(2.3), "linear", "Linear Pattern")
    draw_scatter_panel(slide, Inches(10.0), Inches(1.2), Inches(2.9), Inches(2.3), "nonlinear", "Nonlinear Pattern")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Pearson’s r (Concept)",
        [
            "Pearson’s r measures linear correlation between two variables.",
            "Range: -1 to +1, where sign gives direction.",
            "Closer to ±1 means stronger linear association.",
            "Values near 0 indicate weak or no linear pattern.",
            "Interpret context: strength thresholds vary by field.",
        ],
    )
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.0), Inches(2.2), Inches(5.2), Inches(0.5))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.color.rgb = ACCENT
    for i, txt in enumerate(["-1", "-0.5", "0", "+0.5", "+1"]):
        x = Inches(7.0 + i * 1.3)
        tick = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, x, Inches(2.8), Inches(0.22), Inches(0.22))
        tick.fill.solid()
        tick.fill.fore_color.rgb = BLUE
        tick.line.fill.background()
        t = slide.shapes.add_textbox(x - Inches(0.1), Inches(3.05), Inches(0.5), Inches(0.3))
        p = t.text_frame.paragraphs[0]
        p.text = txt
        p.font.size = Pt(12)
        p.font.color.rgb = BLUE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Pearson’s r (Computation & Interpretation)",
        [
            "Formula idea: compare joint variation to individual spread.",
            "Step 1: compute means of X and Y.",
            "Step 2: calculate paired deviations and products.",
            "Step 3: divide by spread to get standardized r.",
            "Avoid mistakes: outliers, nonlinearity, and causation claims.",
        ],
    )
    formula = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(1.9), Inches(5.5), Inches(1.1))
    formula.fill.solid()
    formula.fill.fore_color.rgb = LIGHT_BLUE
    formula.line.color.rgb = ACCENT
    f = formula.text_frame.paragraphs[0]
    f.text = "r = Σ[(x - x̄)(y - ȳ)] / √(Σ(x - x̄)² Σ(y - ȳ)²)"
    f.font.size = Pt(22)
    f.font.bold = True
    f.font.color.rgb = BLUE
    flow1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(7.1), Inches(3.4), Inches(1.6), Inches(0.9))
    flow2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(8.8), Inches(3.4), Inches(1.6), Inches(0.9))
    flow3 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(10.5), Inches(3.4), Inches(1.6), Inches(0.9))
    for box, text in [(flow1, "Mean"), (flow2, "Covary"), (flow3, "Scale")]:
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = ACCENT
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.font.size = Pt(14)
        p.alignment = 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Hypothesis Testing in Correlation",
        [
            "Hypothesis testing checks if observed r is statistically real.",
            "H₀: r = 0 (no linear relationship in population).",
            "H₁: r ≠ 0 (there is a linear relationship).",
            "Choose significance level α (commonly 0.05).",
            "Goal: decide whether evidence against H₀ is strong enough.",
        ],
    )
    box1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.0), Inches(2.4), Inches(1.0))
    box2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.0), Inches(2.0), Inches(2.4), Inches(1.0))
    for box, text in [(box1, "H₀: r = 0"), (box2, "H₁: r ≠ 0")]:
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = ACCENT
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BLUE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Testing Significance of Correlation",
        [
            "Compute test statistic from sample r and sample size n.",
            "Compare p-value or |t| with critical threshold.",
            "Reject H₀ when evidence is strong at chosen α.",
            "Fail to reject H₀ when evidence is insufficient.",
            "Interpret in context: practical vs statistical importance.",
        ],
    )
    graph = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.2), Inches(5.9), Inches(3.8))
    graph.fill.solid()
    graph.fill.fore_color.rgb = RGBColor(248, 251, 255)
    graph.line.color.rgb = ACCENT
    left_tail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(7.25), Inches(3.4), Inches(1.2), Inches(1.0))
    center = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(2.3), Inches(2.5), Inches(1.6))
    right_tail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(11.3), Inches(3.4), Inches(1.2), Inches(1.0))
    for shape, text in [(left_tail, "Reject"), (center, "Fail to Reject"), (right_tail, "Reject")]:
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BLUE
        shape.line.color.rgb = ACCENT
        p = shape.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.alignment = 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "T-Test (Concept and Formula)",
        [
            "The t-test evaluates if sample correlation differs from zero.",
            "Used when testing significance of Pearson’s r.",
            "Depends on r and sample size n.",
            "Degrees of freedom: n - 2.",
            "Large |t| values provide stronger evidence against H₀.",
        ],
    )
    tbox = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.0), Inches(5.7), Inches(1.2))
    tbox.fill.solid()
    tbox.fill.fore_color.rgb = LIGHT_BLUE
    tbox.line.color.rgb = ACCENT
    pt = tbox.text_frame.paragraphs[0]
    pt.text = "t = r √(n - 2) / √(1 - r²)"
    pt.font.size = Pt(30)
    pt.font.bold = True
    pt.font.color.rgb = BLUE
    small = slide.shapes.add_textbox(Inches(7.2), Inches(3.4), Inches(5.2), Inches(1.2))
    small.text = "Variables: r = sample correlation, n = sample size"
    small.text_frame.paragraphs[0].font.size = Pt(16)
    small.text_frame.paragraphs[0].font.color.rgb = DARK

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "T-Test (Application)",
        [
            "Example: r = 0.62, n = 25, α = 0.05 (two-tailed).",
            "Compute t ≈ 3.79 and df = 23.",
            "Critical value ≈ ±2.07 at α = 0.05.",
            "Since |3.79| > 2.07, reject H₀.",
            "Conclusion: correlation is statistically significant.",
        ],
    )
    sample = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.0), Inches(1.9), Inches(5.7), Inches(2.6))
    sample.fill.solid()
    sample.fill.fore_color.rgb = LIGHT_BLUE
    sample.line.color.rgb = ACCENT
    pp = sample.text_frame.paragraphs[0]
    pp.text = "Quick calculation:\n t = 0.62 × √23 / √(1 - 0.3844)\n t ≈ 3.79"
    pp.font.size = Pt(20)
    pp.font.color.rgb = BLUE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Conclusion / Summary",
        [
            "Correlation quantifies direction and strength of relationships.",
            "Pearson’s r is best for linear, continuous variables.",
            "Significance testing supports evidence-based conclusions.",
            "Always check assumptions, outliers, and context.",
            "Applications: education, health, economics, and business.",
        ],
    )
    for i, txt in enumerate(["Analyze", "Test", "Interpret"]):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(7.0 + i * 1.8), Inches(2.2), Inches(1.7), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = ACCENT
        tp = box.text_frame.paragraphs[0]
        tp.text = txt
        tp.font.bold = True
        tp.font.size = Pt(16)
        tp.font.color.rgb = BLUE
        tp.alignment = 1

    prs.save(str(output_path))


def main():
    parser = argparse.ArgumentParser(description="Generate a presentation on the Measure of Correlation.")
    parser.add_argument(
        "-o",
        "--output",
        default="measure_of_correlation.pptx",
        help="Output presentation path (default: measure_of_correlation.pptx)",
    )
    args = parser.parse_args()
    build_presentation(Path(args.output).resolve())
    print(f"Presentation created: {Path(args.output).resolve()}")


if __name__ == "__main__":
    main()
